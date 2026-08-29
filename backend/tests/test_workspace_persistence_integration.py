import asyncio
import io
import os
import uuid
from contextlib import suppress

import httpx
import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from app.config import Settings
from app.main import create_app
from app.persistence.session_repository import SessionRepository
from app.ppt.domain import PptStatus

pytestmark = pytest.mark.integration


def _settings() -> Settings:
    if os.getenv("RUN_WORKSPACE_PERSISTENCE_INTEGRATION", "").strip().lower() not in {"1", "true", "yes"}:
        pytest.skip(
            "set RUN_WORKSPACE_PERSISTENCE_INTEGRATION=1 with real DATABASE_URL and MinIO settings"
        )

    settings = Settings()
    if settings.persistence_mode != "database":
        pytest.fail("PERSISTENCE_MODE=database is required")
    if not settings.database_url.strip():
        pytest.fail("DATABASE_URL is required")
    if not settings.minio_endpoint.strip():
        pytest.fail("MINIO_ENDPOINT is required")
    if not settings.minio_access_key.strip() or not settings.minio_secret_key.strip():
        pytest.fail("MINIO_ACCESS_KEY and MINIO_SECRET_KEY are required")
    if not settings.minio_public_read:
        pytest.fail("MINIO_PUBLIC_READ=true is required for persisted asset URL verification")
    return settings


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _assert_public_object(url: str, expected: bytes) -> None:
    with httpx.Client(timeout=10, trust_env=False) as client:
        response = client.get(url)
    response.raise_for_status()
    assert response.content == expected


def _cleanup(
    app,
    *,
    file_id: str,
    ppt_id: int,
    ppt_object_name: str,
    conversation_id: str,
    ppt_conversation_id: str,
) -> None:
    file_service = app.state.file_service
    if file_id and file_service is not None:
        with suppress(Exception):
            file_service.delete(file_id)

    ppt_repository = app.state.ppt_repository
    if ppt_id and ppt_repository is not None:
        with suppress(Exception):
            ppt_repository.delete(ppt_id)

    object_store = app.state.object_store
    if object_store is not None:
        with suppress(Exception):
            object_store.delete(ppt_object_name)

    database = app.state.database
    if database is not None:
        with database.session_factory() as session:
            repository = SessionRepository(session)
            repository.delete_conversation(conversation_id)
            repository.delete_conversation(ppt_conversation_id)
            session.commit()


def test_workspace_assets_survive_app_restart() -> None:
    settings = _settings()
    token = uuid.uuid4().hex
    conversation_id = f"h8-chat-{token}"
    ppt_conversation_id = f"h8-ppt-{token}"
    ppt_object_name = f"ppt/{ppt_conversation_id}/h8-{token}.pptx"
    file_id = ""
    ppt_id = 0
    ppt_url = ""
    ppt_payload = _pptx_bytes()

    first_app = create_app(settings)
    cleanup_app = first_app
    try:
        with TestClient(first_app) as client:
            ready = client.get("/health/ready")
            assert ready.status_code == 200
            checks = ready.json()["checks"]
            assert checks["database"]["status"] == "ok"
            assert checks["minio"]["status"] == "ok"
            if settings.vector_database_url.strip():
                assert checks["pgvector"]["status"] == "ok"

            handle = asyncio.run(
                first_app.state.memory.begin_turn(
                    conversation_id,
                    "Persist this conversation across restart.",
                    agent_type="websearch",
                )
            )
            asyncio.run(
                first_app.state.memory.finish_turn(
                    handle,
                    question="Persist this conversation across restart.",
                    answer="Persistent answer.",
                    thinking="Persistent thinking.",
                )
            )

            upload = client.post(
                "/file/upload",
                files={"file": ("h8-persistence.txt", b"persistent file content", "text/plain")},
            ).json()
            assert upload["code"] == 200
            file_id = upload["data"]["fileId"]
            file_url = upload["data"]["minioPath"]
            assert file_url
            _assert_public_object(file_url, b"persistent file content")

            repository = first_app.state.ppt_repository
            assert repository is not None
            inst = repository.create_inst(ppt_conversation_id, "Persistent presentation")
            ppt_id = inst.id

            object_store = first_app.state.object_store
            assert object_store is not None
            ppt_url = object_store.upload(
                object_name=ppt_object_name,
                content=ppt_payload,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            repository.update_file_url(ppt_id, ppt_url, PptStatus.SUCCESS)
            _assert_public_object(ppt_url, ppt_payload)

        second_app = create_app(settings)
        cleanup_app = second_app
        with TestClient(second_app) as client:
            session_detail = client.get(f"/session/{conversation_id}").json()
            assert session_detail["code"] == 200
            assert session_detail["data"]["agentType"] == "websearch"
            assert session_detail["data"]["messages"][0]["answer"] == "Persistent answer."

            sessions = client.get("/session/list", params={"pageNum": 1, "pageSize": 100}).json()
            assert sessions["code"] == 200
            assert any(
                item["conversationId"] == conversation_id
                for item in sessions["data"]["records"]
            )

            file_detail = client.get(f"/file/info/{file_id}").json()
            assert file_detail["code"] == 200
            assert file_detail["data"]["fileName"] == "h8-persistence.txt"
            assert file_detail["data"]["status"] == "SUCCESS"
            file_content = client.get(f"/file/content/{file_id}").json()
            assert file_content["data"]["content"] == "persistent file content"
            _assert_public_object(file_detail["data"]["minioPath"], b"persistent file content")

            ppt_detail = client.get(f"/ppt/{ppt_id}").json()
            assert ppt_detail["code"] == 200
            assert ppt_detail["data"]["conversationId"] == ppt_conversation_id
            assert ppt_detail["data"]["status"] == "SUCCESS"
            assert ppt_detail["data"]["fileUrl"] == ppt_url
            presentations = client.get("/ppt/list").json()
            assert presentations["code"] == 200
            assert any(item["id"] == ppt_id for item in presentations["data"]["presentations"])
            _assert_public_object(ppt_url, ppt_payload)

            assert client.delete(f"/file/{file_id}").json()["code"] == 200
            file_id = ""
            assert client.delete(f"/ppt/{ppt_id}").json()["code"] == 200
            ppt_id = 0
            assert client.delete(f"/session/{conversation_id}").json()["code"] == 200
    finally:
        _cleanup(
            cleanup_app,
            file_id=file_id,
            ppt_id=ppt_id,
            ppt_object_name=ppt_object_name,
            conversation_id=conversation_id,
            ppt_conversation_id=ppt_conversation_id,
        )
