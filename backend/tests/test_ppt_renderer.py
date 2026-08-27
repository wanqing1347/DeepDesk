import asyncio
import json
from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.config import Settings
from app.persistence.models import AiPptInst, AiPptTemplate
from app.ppt.providers import PythonPptRenderer


class CapturingObjectStore:
    def __init__(self) -> None:
        self.uploads: list[tuple[str, bytes, str]] = []

    def upload(self, *, object_name: str, content: bytes, content_type: str) -> str:
        self.uploads.append((object_name, content, content_type))
        return f"https://files.example/{object_name}"

    def delete(self, object_name: str) -> None:
        return None


def _template(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.name = "title"
    box.text = "template title"
    presentation.save(path)


def _two_page_template(path: Path) -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    first_title = first.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    first_title.name = "title"
    first_title.text = "first template"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_title = second.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    second_title.name = "title"
    second_title.text = "second template"
    marker = second.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    marker.name = "marker"
    marker.text = "SECOND_TEMPLATE_MARKER"
    presentation.save(path)


def _inst() -> AiPptInst:
    return AiPptInst(id=1, conversation_id="ppt-render", template_code="ai", status="RENDER")


def _template_record(path: Path) -> AiPptTemplate:
    return AiPptTemplate(
        id=1,
        template_code="ai",
        template_name="AI",
        template_schema='{"slides":[]}',
        file_path=str(path),
    )


def test_python_ppt_renderer_executes_bundled_render_script_and_uploads_valid_pptx(tmp_path: Path) -> None:
    async def scenario() -> None:
        template_path = tmp_path / "template.pptx"
        _template(template_path)
        store = CapturingObjectStore()
        renderer = PythonPptRenderer(
            Settings(
                ppt_render_script_path="./resources/python/render_ppt.py",
                ppt_output_dir=str(tmp_path / "output"),
                ppt_schema_env_threshold_chars=1,
                ppt_render_timeout_seconds=30,
            ),
            store,
        )
        schema = json.dumps(
            {
                "slides": [
                    {
                        "pageType": "COVER",
                        "pageDesc": "cover",
                        "templatePageIndex": 1,
                        "data": {
                            "title": {
                                "type": "text",
                                "content": "rendered title",
                                "fontLimit": 30,
                            }
                        },
                    }
                ]
            },
            ensure_ascii=False,
        )

        url = await renderer.render(_inst(), _template_record(template_path), schema)

        assert url.startswith("https://files.example/ppt/ppt-render/")
        assert len(store.uploads) == 1
        object_name, ppt_bytes, content_type = store.uploads[0]
        assert object_name.startswith("ppt/ppt-render/ppt_1_")
        assert content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        rendered = Presentation(BytesIO(ppt_bytes))
        assert len(rendered.slides) == 1
        text = "\n".join(
            shape.text
            for shape in rendered.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert "rendered title" in text
        assert list((tmp_path / "output").glob("*.pptx")) == []

    asyncio.run(scenario())


def test_python_ppt_renderer_accepts_page_index_alias_for_template_selection(tmp_path: Path) -> None:
    async def scenario() -> None:
        template_path = tmp_path / "two-pages.pptx"
        _two_page_template(template_path)
        store = CapturingObjectStore()
        renderer = PythonPptRenderer(
            Settings(
                ppt_render_script_path="./resources/python/render_ppt.py",
                ppt_output_dir=str(tmp_path / "output"),
                ppt_render_timeout_seconds=30,
            ),
            store,
        )
        schema = json.dumps(
            {
                "slides": [
                    {
                        "pageType": "CONTENT",
                        "pageDesc": "content",
                        "pageIndex": 2,
                        "data": {
                            "title": {
                                "type": "text",
                                "content": "rendered from pageIndex",
                                "fontLimit": 30,
                            }
                        },
                    }
                ]
            },
            ensure_ascii=False,
        )

        await renderer.render(_inst(), _template_record(template_path), schema)

        assert len(store.uploads) == 1
        rendered = Presentation(BytesIO(store.uploads[0][1]))
        assert len(rendered.slides) == 1
        text = "\n".join(
            shape.text
            for shape in rendered.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert "rendered from pageIndex" in text
        assert "SECOND_TEMPLATE_MARKER" in text

    asyncio.run(scenario())


def test_python_ppt_renderer_times_out_and_terminates_child_process(tmp_path: Path) -> None:
    async def scenario() -> None:
        template_path = tmp_path / "template.pptx"
        _template(template_path)
        sleeper = tmp_path / "sleep_renderer.py"
        sleeper.write_text("import time\ntime.sleep(30)\n", encoding="utf-8")
        renderer = PythonPptRenderer(
            Settings(
                ppt_render_script_path=str(sleeper),
                ppt_output_dir=str(tmp_path / "output"),
                ppt_render_timeout_seconds=1,
            ),
            CapturingObjectStore(),
        )

        with pytest.raises(RuntimeError, match="渲染超时"):
            await renderer.render(_inst(), _template_record(template_path), '{"slides":[]}')

    asyncio.run(scenario())


def test_python_ppt_renderer_cancellation_propagates_and_stops_child_process(tmp_path: Path) -> None:
    async def scenario() -> None:
        template_path = tmp_path / "template.pptx"
        _template(template_path)
        sleeper = tmp_path / "sleep_renderer.py"
        sleeper.write_text("import time\ntime.sleep(30)\n", encoding="utf-8")
        renderer = PythonPptRenderer(
            Settings(
                ppt_render_script_path=str(sleeper),
                ppt_output_dir=str(tmp_path / "output"),
                ppt_render_timeout_seconds=30,
            ),
            CapturingObjectStore(),
        )

        task = asyncio.create_task(renderer.render(_inst(), _template_record(template_path), '{"slides":[]}'))
        await asyncio.sleep(0.2)
        task.cancel()
        with pytest.raises(asyncio.CancelledError):
            await asyncio.wait_for(task, timeout=3)

    asyncio.run(scenario())
